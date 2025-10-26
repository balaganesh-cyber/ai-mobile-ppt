#!/usr/bin/env python3
"""
Generate PPTX and upload to Dropbox, then print a shareable link.

Usage:
  1) pip install python-pptx pillow requests
  2) Set env var DROPBOX_TOKEN or paste token when prompted.
     Create token: https://www.dropbox.com/developers/apps (create app -> scoped access -> files.content.write + sharing)
  3) python generate_and_upload_to_dropbox.py
"""
import os
import io
import base64
import getpass
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUT_FILE = "AI_in_Mobile_Phones_Droplet.pptx"
IMG_DIR = "ppt_images"
os.makedirs(IMG_DIR, exist_ok=True)

IMAGE_QUERIES = {
    "title": "smartphone,technology",
    "agenda": "icons,agenda",
    "what_is": "mobile,ai",
    "on_vs_cloud": "cloud,phone,diagram",
    "hardware": "chip,npu",
    "software": "tensorflow,code",
    "photo_voice": "camera,voice,wave",
    "personal_security": "face,unlock,ar",
    "workflow": "pipeline,workflow",
    "challenges": "warning,ethics",
    "future": "futuristic,phone",
    "references": "books,links"
}

def download_placeholder(name, query, max_w=1600, max_h=900):
    url = f"https://source.unsplash.com/featured/{max_w}x{max_h}/?{query}"
    try:
        r = requests.get(url, timeout=15)
        r.raise_for_status()
        path = os.path.join(IMG_DIR, f"{name}.jpg")
        with open(path, "wb") as f:
            f.write(r.content)
        return path
    except Exception as e:
        print(f"[!] image download failed for {name}: {e}")
        return None

slides = [
    {"title":"AI in Mobile Phones","bullets":["How on-device intelligence is changing mobile experiences"],"notes":"Introduce yourself and objectives.","image":"title"},
    {"title":"Agenda","bullets":["What is mobile AI?","On-device vs cloud","Hardware & software","Use cases","Developer workflow","Challenges & ethics","Future trends","Appendix & references"],"notes":"Walk through the agenda.","image":"agenda"},
    {"title":"What is Mobile AI?","bullets":["ML models running on/for mobile devices","Goals: personalization, performance, privacy, context awareness","Constraints: power, memory, latency"],"notes":"Define mobile AI and constraints.","image":"what_is"},
    {"title":"On-device vs Cloud AI","bullets":["On-device: low latency, privacy, offline capable","Cloud: heavy compute, aggregated analytics, higher latency","Hybrid: split execution, federated learning"],"notes":"Examples: wake word on-device vs heavy NLP in cloud.","image":"on_vs_cloud"},
    {"title":"Enabling Hardware","bullets":["NPUs, GPUs, DSPs","Memory & storage; sensors (camera, mic, IMU)","Sensor fusion & accelerators"],"notes":"Explain role of NPUs and sensors.","image":"hardware"},
    {"title":"Enabling Software & Tooling","bullets":["Training: TensorFlow, PyTorch","Edge runtimes: TFLite, Core ML, ONNX","Optimizations: quantization, pruning"],"notes":"Mention converters and quantization.","image":"software"},
    {"title":"Key Use Cases — Photography & Voice","bullets":["Computational photography: HDR, Night Mode","Voice: wake words, on-device speech recognition"],"notes":"Vendor examples like Pixel and Apple.","image":"photo_voice"},
    {"title":"Key Use Cases — Personalization, Security & AR","bullets":["Predictive text & recommendations","Face unlock & anti-spoofing","AR overlays and scene understanding"],"notes":"Highlight privacy benefits.","image":"personal_security"},
    {"title":"Developer Workflow","bullets":["Train in cloud → optimize → convert to TFLite/CoreML/ONNX","Integrate, test & profile on device, monitor & update"],"notes":"Profile for latency, memory and power.","image":"workflow"},
    {"title":"Challenges & Ethics","bullets":["Battery & thermal limits","Model updates & compatibility","Privacy, bias & adversarial attacks"],"notes":"Mitigations: differential privacy, audits.","image":"challenges"},
    {"title":"Future Trends","bullets":["Stronger on-device models as NPUs scale","Continuous personalized on-device learning","Multimodal AI and privacy-first architectures"],"notes":"Suggest student project ideas.","image":"future"},
    {"title":"Appendix — References & Glossary","bullets":["TensorFlow Lite: https://www.tensorflow.org/lite","Core ML: https://developer.apple.com/documentation/coreml","ONNX Runtime: https://onnxruntime.ai","Paper: Federated Learning (McMahan et al.)"],"notes":"References and glossary.","image":"references"}
]

def build_presentation(image_paths):
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    TITLE_RGB = RGBColor(0x00, 0x6D, 0x75)
    for s in slides:
        slide = prs.slides.add_slide(blank_layout)
        # background approximation
        try:
            bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(0x00, 0x8A, 0x92)
            bg.line.fill.background()
        except Exception:
            pass
        # title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1.2))
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = s["title"]
        p.font.size = Pt(34)
        p.font.bold = True
        p.font.color.rgb = TITLE_RGB
        # bullets
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(5.3), Inches(4.6))
        btf = body.text_frame
        first = True
        for b in s["bullets"]:
            if first:
                p = btf.paragraphs[0]; p.text = b; first = False
            else:
                p = btf.add_paragraph(); p.text = b
            p.font.size = Pt(18); p.font.color.rgb = RGBColor(255,255,255)
        # image
        ip = image_paths.get(s["image"])
        if ip and os.path.exists(ip):
            try:
                slide.shapes.add_picture(ip, Inches(6.2), Inches(1.6), width=Inches(3.4))
            except Exception as e:
                print(f"[!] insert image failed: {e}")
        # footer and notes
        footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.55), Inches(9), Inches(0.4))
        footer.text_frame.text = "AI in Mobile Phones — Droplet theme"
        slide.notes_slide.notes_text_frame.text = s["notes"]
    prs.save(OUT_FILE)
    print(f"[+] Saved presentation to {OUT_FILE}")

def upload_to_dropbox(token, local_path, dest_path="/AI_in_Mobile_Phones_Droplet.pptx"):
    # Upload file
    url = "https://content.dropboxapi.com/2/files/upload"
    headers = {
        "Authorization": f"Bearer {token}",
        "Dropbox-API-Arg": '{"path": "%s", "mode": "overwrite", "autorename": false, "mute": false}' % dest_path,
        "Content-Type": "application/octet-stream"
    }
    with open(local_path, "rb") as f:
        data = f.read()
    r = requests.post(url, headers=headers, data=data)
    if r.status_code not in (200, 201):
        print("[!] Dropbox upload failed:", r.status_code, r.text)
        return None
    print("[+] Uploaded to Dropbox:", dest_path)
    # Create shared link
    share_url = "https://api.dropboxapi.com/2/sharing/create_shared_link_with_settings"
    headers = {"Authorization": f"Bearer {token}", "Content-Type": "application/json"}
    payload = {"path": dest_path, "settings": {"requested_visibility": "public"}}
    r2 = requests.post(share_url, headers=headers, json=payload)
    if r2.status_code in (200,201):
        link = r2.json().get("url")
        # Dropbox returns a dl=0 link; convert to dl=1 for direct download
        if link and "dl=0" in link:
            link = link.replace("?dl=0", "?dl=1")
        print("[+] Shareable link:", link)
        return link
    else:
        # If link exists, try to list and fetch
        print("[!] create_shared_link failed:", r2.status_code, r2.text)
        return None

def main():
    token = os.getenv("DROPBOX_TOKEN")
    if not token:
        token = getpass.getpass("Enter your Dropbox API token (input hidden): ").strip()
    print("Downloading images...")
    image_paths = {}
    for k,q in IMAGE_QUERIES.items():
        image_paths[k] = download_placeholder(k, q)
    print("Building presentation...")
    build_presentation(image_paths)
    print("Uploading to Dropbox...")
    link = upload_to_dropbox(token, OUT_FILE)
    if link:
        print("\nDone. Use the link above to download the PPTX.")
    else:
        print("\nUpload completed with issues. Check errors above.")

if __name__ == "__main__":
    main()